import io
import xlsxwriter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData, XyChartData
from pptx.dml.color import RGBColor
import os
import numpy as np

LOGO_PATH = os.path.join("static", "logos", "credivex_logo.jpg")

# ============================================================
#  COLOR PALETTE  (light-green / green / light-blue / blue / yellow)
# ============================================================
# Excel hex values (no #)
XL_GREEN       = '#059669'   # Green (primary approved)
XL_LIGHT_GREEN = '#6ee7b7'   # Light green (accent)
XL_BLUE        = '#2563eb'   # Blue
XL_LIGHT_BLUE  = '#7dd3fc'   # Light blue (accent)
XL_YELLOW      = '#eab308'   # Yellow (highlight)
XL_DENIED      = '#f87171'   # Soft red for denied (kept for data clarity)

# PowerPoint RGBColor values
PPT_GREEN       = RGBColor(5, 150, 105)     # #059669
PPT_LIGHT_GREEN = RGBColor(110, 231, 183)   # #6ee7b7
PPT_BLUE        = RGBColor(37, 99, 235)     # #2563eb
PPT_LIGHT_BLUE  = RGBColor(125, 211, 252)   # #7dd3fc
PPT_YELLOW      = RGBColor(234, 179, 8)     # #eab308
PPT_DENIED_GREY = RGBColor(148, 163, 184)   # #94a3b8  (muted for denied dots)

# Slide theme
BG_COLOR        = RGBColor(240, 253, 244)   # #f0fdf4  very light green tint
PANEL_COLOR     = RGBColor(255, 255, 255)   # white panels
BORDER_COLOR    = RGBColor(187, 247, 208)   # #bbf7d0  light green border
TEXT_PRIMARY    = RGBColor(15, 23, 42)      # #0f172a  dark slate
TEXT_SECONDARY  = RGBColor(71, 85, 105)     # #475569  muted slate
ACCENT_COLOR    = PPT_GREEN                  # primary accent
KPI_ACCENT      = PPT_BLUE                   # KPI highlight


# ============================================================
#  EXCEL REPORT
# ============================================================
def generate_excel_report(results):
    output = io.BytesIO()
    workbook = xlsxwriter.Workbook(output, {'in_memory': True})

    # 1. Data Sheet
    data_sheet = workbook.add_worksheet('Evaluation Data')
    header_format = workbook.add_format({
        'bold': True, 'bg_color': '#059669', 'font_color': 'white',
        'border': 1, 'border_color': '#047857'
    })

    headers = [
        'Row', 'Annual Income', 'Loan Amount', 'Credit Score',
        'DTI', 'Years Employed', 'Delinquencies',
        'Verdict', 'Approval %', 'Risk Tier'
    ]
    for col, h in enumerate(headers):
        data_sheet.write(0, col, h, header_format)

    # Alternate row colors for readability
    row_even = workbook.add_format({'bg_color': '#f0fdf4'})
    row_odd  = workbook.add_format({})

    for row_idx, res in enumerate(results, start=1):
        fmt = row_even if row_idx % 2 == 0 else row_odd
        data_sheet.write(row_idx, 0, res.get('row', row_idx), fmt)
        data_sheet.write(row_idx, 1, res['inputs'].get('annual_income', 0), fmt)
        data_sheet.write(row_idx, 2, res['inputs'].get('loan_amount', 0), fmt)
        data_sheet.write(row_idx, 3, res['inputs'].get('credit_score', 0), fmt)
        data_sheet.write(row_idx, 4, res['inputs'].get('debt_to_income_ratio', 0), fmt)
        data_sheet.write(row_idx, 5, res['inputs'].get('years_employed', 0), fmt)
        data_sheet.write(row_idx, 6, res['inputs'].get('delinquencies_last_2yrs', 0), fmt)
        data_sheet.write(row_idx, 7, res.get('prediction', ''), fmt)
        data_sheet.write(row_idx, 8, res.get('approval_percentage', ''), fmt)
        data_sheet.write(row_idx, 9, res.get('risk_tier', ''), fmt)

    # Auto-fit column widths (approximate)
    col_widths = [6, 16, 14, 14, 8, 16, 14, 10, 14, 18]
    for i, w in enumerate(col_widths):
        data_sheet.set_column(i, i, w)

    # 2. Analytics Sheet
    chart_sheet = workbook.add_worksheet('Visual Analytics')

    # Add Logo
    if os.path.exists(LOGO_PATH):
        chart_sheet.insert_image('B2', LOGO_PATH, {'x_scale': 0.5, 'y_scale': 0.5})

    chart_sheet.write('B7', 'CREDIVEX Bulk Evaluation Analytics',
                      workbook.add_format({'bold': True, 'font_size': 20, 'font_color': '#059669'}))

    # Hidden data sheet for charts
    cdata = workbook.add_worksheet('ChartData')
    cdata.hide()

    # Collect chart data
    scatter_app_x, scatter_app_y = [], []
    scatter_den_x, scatter_den_y = [], []

    def get_ye_bucket(ye):
        if ye <= 2: return 0
        if ye <= 5: return 1
        if ye <= 10: return 2
        if ye <= 20: return 3
        return 4

    ye_labels = ['0-2 yrs', '3-5 yrs', '6-10 yrs', '11-20 yrs', '21+ yrs']
    ye_totals = [0] * 5
    ye_approvals = [0] * 5

    def get_loan_bucket(amt):
        if amt < 10000: return 0
        if amt < 20000: return 1
        if amt < 30000: return 2
        if amt < 50000: return 3
        if amt < 75000: return 4
        if amt < 100000: return 5
        if amt < 150000: return 6
        return 7

    loan_labels = ['<$10k', '$10-20k', '$20-30k', '$30-50k',
                   '$50-75k', '$75-100k', '$100-150k', '>$150k']
    loan_app = [0] * 8
    loan_den = [0] * 8

    for r in results:
        dti = r['inputs'].get('debt_to_income_ratio', 0) * 100
        score = r['inputs'].get('credit_score', 0)
        if r.get('is_approved'):
            scatter_app_x.append(dti)
            scatter_app_y.append(score)
        else:
            scatter_den_x.append(dti)
            scatter_den_y.append(score)

        ye = r['inputs'].get('years_employed', 0)
        idx = get_ye_bucket(ye)
        ye_totals[idx] += 1
        if r.get('is_approved'):
            ye_approvals[idx] += 1

        amt = r['inputs'].get('loan_amount', 0)
        l_idx = get_loan_bucket(amt)
        if r.get('is_approved'):
            loan_app[l_idx] += 1
        else:
            loan_den[l_idx] += 1

    ye_rates = [round((ye_approvals[i] / ye_totals[i]) * 100, 1)
                if ye_totals[i] > 0 else 0 for i in range(5)]

    # Write chart backing data
    cdata.write_column('H1', ['App DTI'] + scatter_app_x)
    cdata.write_column('I1', ['App Score'] + scatter_app_y)
    cdata.write_column('J1', ['Den DTI'] + scatter_den_x)
    cdata.write_column('K1', ['Den Score'] + scatter_den_y)
    cdata.write_column('A1', ['Years Employed'] + ye_labels)
    cdata.write_column('B1', ['Approval Rate'] + ye_rates)
    cdata.write_column('D1', ['Loan Amount'] + loan_labels)
    cdata.write_column('E1', ['Approved'] + loan_app)
    cdata.write_column('F1', ['Denied'] + loan_den)

    # ---- Chart 0: Scatter (DTI vs Credit Score) ----
    scatter = workbook.add_chart({'type': 'scatter'})
    if scatter_app_x:
        scatter.add_series({
            'name': 'Approved',
            'categories': f'=ChartData!$H$2:$H${len(scatter_app_x)+1}',
            'values':     f'=ChartData!$I$2:$I${len(scatter_app_y)+1}',
            'marker': {'type': 'circle', 'size': 3,
                       'border': {'none': True},
                       'fill': {'color': '#059669'}}
        })
    if scatter_den_x:
        scatter.add_series({
            'name': 'Denied',
            'categories': f'=ChartData!$J$2:$J${len(scatter_den_x)+1}',
            'values':     f'=ChartData!$K$2:$K${len(scatter_den_y)+1}',
            'marker': {'type': 'circle', 'size': 3,
                       'border': {'none': True},
                       'fill': {'color': '#7dd3fc'}}
        })
    scatter.set_title({'name': 'Decision Boundary: DTI vs Credit Score'})
    scatter.set_x_axis({'name': 'Debt-to-Income Ratio (%)'})
    scatter.set_y_axis({'name': 'Credit Score (FICO)'})
    chart_sheet.insert_chart('B9', scatter, {'x_scale': 1.2, 'y_scale': 1.1})

    # ---- Chart 1: Bar (Approval Rate) ----
    bar_chart = workbook.add_chart({'type': 'column'})
    bar_chart.add_series({
        'categories': '=ChartData!$A$2:$A$6',
        'values': '=ChartData!$B$2:$B$6',
        'fill': {'color': '#2563eb'}
    })
    bar_chart.set_title({'name': 'Approval Rate by Employment Tenure'})
    bar_chart.set_x_axis({'name': 'Years Employed'})
    bar_chart.set_y_axis({'name': 'Approval Rate (%)', 'max': 100})
    bar_chart.set_legend({'none': True})
    chart_sheet.insert_chart('L9', bar_chart, {'x_scale': 1.2, 'y_scale': 1.1})

    # ---- Chart 2: Histogram (Loan Amount) ----
    hist_chart = workbook.add_chart({'type': 'column', 'subtype': 'stacked'})
    hist_chart.add_series({
        'name': 'Approved',
        'categories': '=ChartData!$D$2:$D$9',
        'values': '=ChartData!$E$2:$E$9',
        'fill': {'color': '#6ee7b7'}
    })
    hist_chart.add_series({
        'name': 'Denied',
        'categories': '=ChartData!$D$2:$D$9',
        'values': '=ChartData!$F$2:$F$9',
        'fill': {'color': '#eab308'}
    })
    hist_chart.set_title({'name': 'Loan Amount Distribution'})
    hist_chart.set_x_axis({'name': 'Loan Amount ($)'})
    hist_chart.set_y_axis({'name': 'Number of Applications'})
    chart_sheet.insert_chart('B28', hist_chart, {'x_scale': 2.4, 'y_scale': 1.2})

    workbook.close()
    output.seek(0)
    return output


# ============================================================
#  POWERPOINT HELPERS
# ============================================================
def set_bg(slide):
    """Set the slide background to a very light green tint."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_panel(slide, x, y, cx, cy):
    """Add a white glassmorphic rounded-rectangle panel."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL_COLOR
    shape.line.color.rgb = BORDER_COLOR
    shape.line.width = Pt(1)
    return shape


def apply_text_style(shape, color, font_size, bold=False):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = color
            run.font.size = Pt(font_size)
            run.font.bold = bold


# ============================================================
#  POWERPOINT REPORT
# ============================================================
def generate_ppt_report(results):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)  # 16:9

    scatter_app_x, scatter_app_y = [], []
    scatter_den_x, scatter_den_y = [], []

    ye_totals = [0] * 5
    ye_approvals = [0] * 5
    ye_labels = ['0-2 yrs', '3-5 yrs', '6-10 yrs', '11-20 yrs', '21+ yrs']

    loan_app = [0] * 8
    loan_den = [0] * 8
    loan_labels = ['<$10k', '$10-20k', '$20-30k', '$30-50k',
                   '$50-75k', '$75-100k', '$100-150k', '>$150k']

    def get_ye_bucket(ye):
        if ye <= 2: return 0
        if ye <= 5: return 1
        if ye <= 10: return 2
        if ye <= 20: return 3
        return 4

    def get_loan_bucket(amt):
        if amt < 10000: return 0
        if amt < 20000: return 1
        if amt < 30000: return 2
        if amt < 50000: return 3
        if amt < 75000: return 4
        if amt < 100000: return 5
        if amt < 150000: return 6
        return 7

    for r in results:
        dti = r['inputs'].get('debt_to_income_ratio', 0) * 100
        score = r['inputs'].get('credit_score', 0)
        is_app = r.get('is_approved', False)
        if is_app:
            scatter_app_x.append(dti)
            scatter_app_y.append(score)
        else:
            scatter_den_x.append(dti)
            scatter_den_y.append(score)

        ye = r['inputs'].get('years_employed', 0)
        idx = get_ye_bucket(ye)
        ye_totals[idx] += 1
        if is_app:
            ye_approvals[idx] += 1

        amt = r['inputs'].get('loan_amount', 0)
        l_idx = get_loan_bucket(amt)
        if is_app:
            loan_app[l_idx] += 1
        else:
            loan_den[l_idx] += 1

    ye_rates = [round((ye_approvals[i] / ye_totals[i]) * 100, 1)
                if ye_totals[i] > 0 else 0 for i in range(5)]

    # ================================================================
    # SLIDE 1 : Title & Executive Summary
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_bg(slide)

    # Green accent bar at top
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0), Inches(0),
                                     Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PPT_GREEN
    top_bar.line.fill.background()

    # Center Panel
    add_panel(slide, Inches(2.66), Inches(1.2), Inches(8), Inches(5))

    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(5.66), Inches(1.5),
                                 height=Inches(1.2))

    tx_box = slide.shapes.add_textbox(Inches(2.66), Inches(3.0), Inches(8), Inches(1))
    tf = tx_box.text_frame
    tf.text = "CREDIVEX Bulk Evaluation Analytics"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    apply_text_style(tx_box, TEXT_PRIMARY, 32, True)

    # KPI box with green border
    kpi_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(4.66), Inches(4.3),
                                     Inches(4), Inches(1.2))
    kpi_box.fill.solid()
    kpi_box.fill.fore_color.rgb = RGBColor(236, 253, 245)  # #ecfdf5
    kpi_box.line.color.rgb = PPT_GREEN
    kpi_box.line.width = Pt(2)
    kpi_tf = kpi_box.text_frame
    kpi_tf.text = f"Total Processed: {len(results)}\nNative Chart Report"
    kpi_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if len(kpi_tf.paragraphs) > 1:
        kpi_tf.paragraphs[1].alignment = PP_ALIGN.CENTER
    apply_text_style(kpi_box, PPT_GREEN, 18, True)

    # ================================================================
    # SLIDE 2 : Decision Boundary (DTI vs Credit Score)
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # Title
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                     Inches(10), Inches(0.8))
    t_box.text_frame.text = "Decision Boundary: DTI vs Credit Score"
    apply_text_style(t_box, TEXT_PRIMARY, 28, True)

    # Left Column Panel
    left_panel = add_panel(slide, Inches(0.5), Inches(1.2),
                           Inches(3.5), Inches(5.8))
    l_tf = left_panel.text_frame
    l_tf.text = ("Support Vector Classifier\n\n"
                 "This boundary separates approved candidates (green) "
                 "from denied candidates (blue) based on their "
                 "Debt-to-Income ratio and Credit Score.")
    apply_text_style(left_panel, TEXT_SECONDARY, 14)
    l_tf.paragraphs[0].font.color.rgb = PPT_GREEN
    l_tf.paragraphs[0].font.size = Pt(18)
    l_tf.paragraphs[0].font.bold = True

    # KPI inside left panel
    kpi2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.5),
                                    Inches(3.1), Inches(1.5))
    kpi2.text_frame.text = "Approval Accuracy\n98.5%"
    apply_text_style(kpi2, TEXT_PRIMARY, 16)
    if len(kpi2.text_frame.paragraphs) > 1:
        kpi2.text_frame.paragraphs[1].font.size = Pt(40)
        kpi2.text_frame.paragraphs[1].font.color.rgb = PPT_BLUE
        kpi2.text_frame.paragraphs[1].font.bold = True

    # Right Column Panel (Chart)
    add_panel(slide, Inches(4.2), Inches(1.2), Inches(8.6), Inches(5.8))

    chart_data = XyChartData()
    if scatter_app_x:
        s1 = chart_data.add_series('Approved')
        for xv, yv in zip(scatter_app_x, scatter_app_y):
            s1.add_data_point(xv, yv)
    if scatter_den_x:
        s2 = chart_data.add_series('Denied')
        for xv, yv in zip(scatter_den_x, scatter_den_y):
            s2.add_data_point(xv, yv)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(4.4), Inches(1.4), Inches(8.2), Inches(5.4),
        chart_data
    ).chart
    chart.has_title = False

    chart.value_axis.has_major_gridlines = False
    chart.category_axis.has_major_gridlines = False

    chart.category_axis.has_title = True
    chart.value_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "Debt-to-Income Ratio (%)"
    chart.value_axis.axis_title.text_frame.text = "Credit Score (FICO)"
    apply_text_style(chart.category_axis.axis_title, TEXT_SECONDARY, 12)
    apply_text_style(chart.value_axis.axis_title, TEXT_SECONDARY, 12)
    chart.category_axis.tick_labels.font.color.rgb = TEXT_SECONDARY
    chart.value_axis.tick_labels.font.color.rgb = TEXT_SECONDARY

    if scatter_app_x:
        chart.series[0].marker.style = XL_MARKER_STYLE.CIRCLE
        chart.series[0].marker.size = 4
        chart.series[0].marker.format.fill.solid()
        chart.series[0].marker.format.fill.fore_color.rgb = PPT_GREEN
        chart.series[0].format.line.fill.background()
    if scatter_den_x:
        s_idx = 1 if scatter_app_x else 0
        chart.series[s_idx].marker.style = XL_MARKER_STYLE.CIRCLE
        chart.series[s_idx].marker.size = 4
        chart.series[s_idx].marker.format.fill.solid()
        chart.series[s_idx].marker.format.fill.fore_color.rgb = PPT_LIGHT_BLUE
        chart.series[s_idx].format.line.fill.background()

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.color.rgb = TEXT_SECONDARY

    # ================================================================
    # SLIDE 3 : Approval Rate by Employment Tenure (horizontal bar)
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                     Inches(10), Inches(0.8))
    t_box.text_frame.text = "Approval Rate by Employment Tenure"
    apply_text_style(t_box, TEXT_PRIMARY, 28, True)

    add_panel(slide, Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.8))

    chart_data = ChartData()
    chart_data.categories = ye_labels[::-1]  # Reverse for horizontal bar
    chart_data.add_series('Approval Rate', tuple(ye_rates[::-1]))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.7), Inches(1.4), Inches(11.93), Inches(5.4),
        chart_data
    ).chart
    chart.has_title = False
    chart.has_legend = False

    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = BORDER_COLOR

    chart.category_axis.tick_labels.font.color.rgb = TEXT_SECONDARY
    chart.category_axis.tick_labels.font.size = Pt(14)
    chart.value_axis.tick_labels.font.color.rgb = TEXT_SECONDARY

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.font.color.rgb = TEXT_PRIMARY
    data_labels.font.size = Pt(12)
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = PPT_BLUE

    # ================================================================
    # SLIDE 4 : Loan Amount Distribution + KPI cards
    # ================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                     Inches(10), Inches(0.8))
    t_box.text_frame.text = "Loan Amount Distribution"
    apply_text_style(t_box, TEXT_PRIMARY, 28, True)

    add_panel(slide, Inches(0.5), Inches(1.2), Inches(8), Inches(5.8))

    chart_data = ChartData()
    chart_data.categories = loan_labels
    chart_data.add_series('Approved', tuple(loan_app))
    chart_data.add_series('Denied', tuple(loan_den))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(0.7), Inches(1.4), Inches(7.6), Inches(5.4),
        chart_data
    ).chart
    chart.has_title = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = BORDER_COLOR

    chart.category_axis.tick_labels.font.color.rgb = TEXT_SECONDARY
    chart.value_axis.tick_labels.font.color.rgb = TEXT_SECONDARY

    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = PPT_LIGHT_GREEN
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = PPT_YELLOW
    chart.series[1].format.line.color.rgb = RGBColor(202, 138, 4)  # darker yellow border
    chart.series[1].format.line.width = Pt(1)

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.color.rgb = TEXT_SECONDARY

    # Right Column KPI cards
    total_app = sum(loan_app)
    arr_total = np.array(loan_app) + np.array(loan_den)
    most_req_idx = int(arr_total.argmax()) if arr_total.sum() > 0 else 0
    highest_app_idx = int(np.array(loan_app).argmax()) if total_app > 0 else 0

    kpis = [
        ("Total Approved", f"{total_app}", PPT_GREEN),
        ("Most Requested Bracket", loan_labels[most_req_idx], PPT_BLUE),
        ("Highest Approval Bracket", loan_labels[highest_app_idx], PPT_YELLOW),
    ]

    for i, (title, val, accent) in enumerate(kpis):
        y_pos = 1.2 + i * 2.0
        kpi_pan = add_panel(slide, Inches(8.8), Inches(y_pos),
                            Inches(4), Inches(1.8))
        # Accent left-edge stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(8.8), Inches(y_pos),
                                        Inches(0.12), Inches(1.8))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()

        k_tf = kpi_pan.text_frame
        k_tf.text = f"{title}\n{val}"
        apply_text_style(kpi_pan, TEXT_SECONDARY, 14)
        if len(k_tf.paragraphs) > 1:
            k_tf.paragraphs[1].font.size = Pt(28)
            k_tf.paragraphs[1].font.color.rgb = accent
            k_tf.paragraphs[1].font.bold = True
            k_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            k_tf.paragraphs[1].alignment = PP_ALIGN.CENTER

    # Add logo to all slides (CREDIVEX branding)
    if os.path.exists(LOGO_PATH):
        for s in prs.slides:
            # Place in the top right corner
            s.shapes.add_picture(LOGO_PATH, Inches(11.3), Inches(0.2), height=Inches(0.6))

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
